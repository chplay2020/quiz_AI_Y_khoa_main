"""
Document Processing Module
Handles extraction of text from various document formats (PDF, PPTX, DOCX)
"""
import os
import asyncio
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import hashlib
import structlog

# PDF Processing
import pdfplumber
from PyPDF2 import PdfReader

# PowerPoint Processing
from pptx import Presentation
from pptx.util import Inches

# Word Processing
from docx import Document

# Text Processing
import re
from dataclasses import dataclass

logger = structlog.get_logger()


@dataclass
class ExtractedChunk:
    """Represents an extracted text chunk from a document"""
    chunk_id: str
    content: str
    page_number: Optional[int]
    section_title: Optional[str]
    metadata: Dict[str, Any]


@dataclass
class ProcessedDocument:
    """Represents a fully processed document"""
    document_id: str
    filename: str
    file_type: str
    total_pages: int
    chunks: List[ExtractedChunk]
    metadata: Dict[str, Any]


class DocumentProcessor:
    """
    Main document processor class that handles text extraction
    from various document formats with medical content awareness.
    """
    
    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        min_chunk_size: int = 100
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size
        
        # Medical-specific patterns
        self.medical_patterns = {
            'diagnosis': r'\b(chẩn đoán|diagnosis|Dx)\b',
            'treatment': r'\b(điều trị|treatment|Tx|thuốc|medication)\b',
            'symptoms': r'\b(triệu chứng|symptoms|biểu hiện|signs)\b',
            'dosage': r'\b(\d+\s*(mg|g|ml|mcg|IU|units?))\b',
            'lab_values': r'\b(\d+\.?\d*\s*(mmol/L|mg/dL|g/dL|U/L|mEq/L))\b',
        }
    
    async def process_file(
        self,
        file_path: str,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessedDocument:
        """
        Process a document file and extract text chunks.
        
        Args:
            file_path: Path to the document file
            document_id: Unique identifier for the document
            metadata: Additional metadata to attach to chunks
        
        Returns:
            ProcessedDocument with extracted chunks
        """
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower()
        
        logger.info("Processing document", file=str(file_path), type=file_extension)
        
        if file_extension == '.pdf':
            return await self._process_pdf(file_path, document_id, metadata)
        elif file_extension in ['.pptx', '.ppt']:
            return await self._process_pptx(file_path, document_id, metadata)
        elif file_extension in ['.docx', '.doc']:
            return await self._process_docx(file_path, document_id, metadata)
        elif file_extension == '.txt':
            return await self._process_txt(file_path, document_id, metadata)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    async def _process_pdf(
        self,
        file_path: Path,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessedDocument:
        """Extract text from PDF documents"""
        chunks = []
        total_pages = 0
        
        try:
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                full_text = []
                page_texts = []
                
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    
                    # Extract tables if present
                    tables = page.extract_tables()
                    table_text = self._format_tables(tables) if tables else ""
                    
                    page_content = f"{text}\n{table_text}".strip()
                    page_texts.append((page_num, page_content))
                    full_text.append(page_content)
                
                # Create chunks with page awareness
                chunks = self._create_chunks_with_pages(
                    page_texts, document_id, metadata or {}
                )
                
        except Exception as e:
            logger.error("PDF processing error", error=str(e))
            raise
        
        return ProcessedDocument(
            document_id=document_id,
            filename=file_path.name,
            file_type="pdf",
            total_pages=total_pages,
            chunks=chunks,
            metadata=metadata or {}
        )
    
    async def _process_pptx(
        self,
        file_path: Path,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessedDocument:
        """Extract text from PowerPoint presentations"""
        chunks = []
        
        try:
            prs = Presentation(str(file_path))
            total_pages = len(prs.slides)
            slide_texts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_title = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.is_placeholder and shape.placeholder_format.type == 1:
                            slide_title = shape.text
                        slide_content.append(shape.text)
                    
                    # Handle tables in slides
                    if shape.has_table:
                        table_text = self._extract_pptx_table(shape.table)
                        slide_content.append(table_text)
                
                full_slide_text = "\n".join(slide_content)
                slide_texts.append((slide_num, slide_title, full_slide_text))
            
            # Create chunks from slides
            chunks = self._create_chunks_from_slides(
                slide_texts, document_id, metadata or {}
            )
            
        except Exception as e:
            logger.error("PPTX processing error", error=str(e))
            raise
        
        return ProcessedDocument(
            document_id=document_id,
            filename=file_path.name,
            file_type="pptx",
            total_pages=total_pages,
            chunks=chunks,
            metadata=metadata or {}
        )
    
    async def _process_docx(
        self,
        file_path: Path,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessedDocument:
        """Extract text from Word documents"""
        chunks = []
        
        try:
            doc = Document(str(file_path))
            paragraphs = []
            current_section = ""
            
            for para in doc.paragraphs:
                # Detect section headers
                if para.style.name.startswith('Heading'):
                    current_section = para.text
                
                if para.text.strip():
                    paragraphs.append({
                        'text': para.text,
                        'section': current_section,
                        'style': para.style.name
                    })
            
            # Extract tables
            for table in doc.tables:
                table_text = self._extract_docx_table(table)
                paragraphs.append({
                    'text': table_text,
                    'section': current_section,
                    'style': 'Table'
                })
            
            # Create chunks
            chunks = self._create_chunks_from_paragraphs(
                paragraphs, document_id, metadata or {}
            )
            
        except Exception as e:
            logger.error("DOCX processing error", error=str(e))
            raise
        
        return ProcessedDocument(
            document_id=document_id,
            filename=file_path.name,
            file_type="docx",
            total_pages=len(chunks),  # Approximate
            chunks=chunks,
            metadata=metadata or {}
        )
    
    async def _process_txt(
        self,
        file_path: Path,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> ProcessedDocument:
        """Extract text from plain text files"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        chunks = self._create_simple_chunks(content, document_id, metadata or {})
        
        return ProcessedDocument(
            document_id=document_id,
            filename=file_path.name,
            file_type="txt",
            total_pages=1,
            chunks=chunks,
            metadata=metadata or {}
        )
    
    def _format_tables(self, tables: List[List[List[str]]]) -> str:
        """Format extracted tables as text"""
        formatted = []
        for table in tables:
            if table:
                for row in table:
                    formatted.append(" | ".join(str(cell) if cell else "" for cell in row))
                formatted.append("")
        return "\n".join(formatted)
    
    def _extract_pptx_table(self, table) -> str:
        """Extract text from a PowerPoint table"""
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
    
    def _extract_docx_table(self, table) -> str:
        """Extract text from a Word table"""
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)
    
    def _create_chunks_with_pages(
        self,
        page_texts: List[Tuple[int, str]],
        document_id: str,
        metadata: Dict[str, Any]
    ) -> List[ExtractedChunk]:
        """Create chunks while preserving page information"""
        chunks = []
        current_chunk = ""
        current_pages = []
        
        for page_num, text in page_texts:
            # Clean the text
            text = self._clean_text(text)
            
            if not text:
                continue
            
            # Split into sentences for better chunking
            sentences = self._split_into_sentences(text)
            
            for sentence in sentences:
                if len(current_chunk) + len(sentence) <= self.chunk_size:
                    current_chunk += sentence + " "
                    if page_num not in current_pages:
                        current_pages.append(page_num)
                else:
                    # Save current chunk if it meets minimum size
                    if len(current_chunk) >= self.min_chunk_size:
                        chunk_id = self._generate_chunk_id(document_id, len(chunks))
                        chunks.append(ExtractedChunk(
                            chunk_id=chunk_id,
                            content=current_chunk.strip(),
                            page_number=current_pages[0] if current_pages else None,
                            section_title=None,
                            metadata={
                                **metadata,
                                'pages': current_pages.copy(),
                                'medical_entities': self._extract_medical_entities(current_chunk)
                            }
                        ))
                    
                    # Start new chunk with overlap
                    overlap_text = current_chunk[-self.chunk_overlap:] if len(current_chunk) > self.chunk_overlap else ""
                    current_chunk = overlap_text + sentence + " "
                    current_pages = [page_num]
        
        # Don't forget the last chunk
        if current_chunk and len(current_chunk) >= self.min_chunk_size:
            chunk_id = self._generate_chunk_id(document_id, len(chunks))
            chunks.append(ExtractedChunk(
                chunk_id=chunk_id,
                content=current_chunk.strip(),
                page_number=current_pages[0] if current_pages else None,
                section_title=None,
                metadata={
                    **metadata,
                    'pages': current_pages,
                    'medical_entities': self._extract_medical_entities(current_chunk)
                }
            ))
        
        return chunks
    
    def _create_chunks_from_slides(
        self,
        slide_texts: List[Tuple[int, str, str]],
        document_id: str,
        metadata: Dict[str, Any]
    ) -> List[ExtractedChunk]:
        """Create chunks from PowerPoint slides"""
        chunks = []
        
        for slide_num, title, content in slide_texts:
            content = self._clean_text(content)
            
            if not content:
                continue
            
            # For slides, we might want to keep each slide as a chunk
            # or merge small slides together
            if len(content) >= self.min_chunk_size:
                chunk_id = self._generate_chunk_id(document_id, len(chunks))
                chunks.append(ExtractedChunk(
                    chunk_id=chunk_id,
                    content=content,
                    page_number=slide_num,
                    section_title=title,
                    metadata={
                        **metadata,
                        'slide_title': title,
                        'medical_entities': self._extract_medical_entities(content)
                    }
                ))
        
        return chunks
    
    def _create_chunks_from_paragraphs(
        self,
        paragraphs: List[Dict[str, str]],
        document_id: str,
        metadata: Dict[str, Any]
    ) -> List[ExtractedChunk]:
        """Create chunks from document paragraphs"""
        chunks = []
        current_chunk = ""
        current_section = ""
        
        for para in paragraphs:
            text = self._clean_text(para['text'])
            section = para.get('section', '')
            
            if not text:
                continue
            
            if len(current_chunk) + len(text) <= self.chunk_size:
                current_chunk += text + "\n"
                current_section = section or current_section
            else:
                if len(current_chunk) >= self.min_chunk_size:
                    chunk_id = self._generate_chunk_id(document_id, len(chunks))
                    chunks.append(ExtractedChunk(
                        chunk_id=chunk_id,
                        content=current_chunk.strip(),
                        page_number=None,
                        section_title=current_section,
                        metadata={
                            **metadata,
                            'medical_entities': self._extract_medical_entities(current_chunk)
                        }
                    ))
                
                overlap_text = current_chunk[-self.chunk_overlap:] if len(current_chunk) > self.chunk_overlap else ""
                current_chunk = overlap_text + text + "\n"
                current_section = section
        
        if current_chunk and len(current_chunk) >= self.min_chunk_size:
            chunk_id = self._generate_chunk_id(document_id, len(chunks))
            chunks.append(ExtractedChunk(
                chunk_id=chunk_id,
                content=current_chunk.strip(),
                page_number=None,
                section_title=current_section,
                metadata={
                    **metadata,
                    'medical_entities': self._extract_medical_entities(current_chunk)
                }
            ))
        
        return chunks
    
    def _create_simple_chunks(
        self,
        text: str,
        document_id: str,
        metadata: Dict[str, Any]
    ) -> List[ExtractedChunk]:
        """Create simple chunks from plain text"""
        text = self._clean_text(text)
        chunks = []
        
        for i in range(0, len(text), self.chunk_size - self.chunk_overlap):
            chunk_text = text[i:i + self.chunk_size]
            
            if len(chunk_text) >= self.min_chunk_size:
                chunk_id = self._generate_chunk_id(document_id, len(chunks))
                chunks.append(ExtractedChunk(
                    chunk_id=chunk_id,
                    content=chunk_text,
                    page_number=None,
                    section_title=None,
                    metadata={
                        **metadata,
                        'medical_entities': self._extract_medical_entities(chunk_text)
                    }
                ))
        
        return chunks
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove common artifacts
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]', '', text)
        
        # Fix common encoding issues
        text = text.replace('\ufeff', '').replace('\u2019', "'")
        
        return text.strip()
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences for better chunking"""
        # Vietnamese and English sentence splitting
        sentence_endings = r'(?<=[.!?。])\s+'
        sentences = re.split(sentence_endings, text)
        return [s.strip() for s in sentences if s.strip()]
    
    def _generate_chunk_id(self, document_id: str, chunk_index: int) -> str:
        """Generate a unique chunk ID"""
        content = f"{document_id}_{chunk_index}"
        return hashlib.md5(content.encode()).hexdigest()[:16]
    
    def _extract_medical_entities(self, text: str) -> Dict[str, List[str]]:
        """Extract medical entities from text"""
        entities = {}
        
        for entity_type, pattern in self.medical_patterns.items():
            matches = re.findall(pattern, text, re.IGNORECASE)
            if matches:
                entities[entity_type] = list(set(matches))
        
        return entities
